# utils/file_parser.py

import os
from typing import Union
import docx2txt
import pptx
import fitz  # PyMuPDF

def extract_text_from_file(file_path: str) -> str:
    """
    Extracts readable text content from PDF, DOCX, or PPTX files.
    """
    if file_path.lower().endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type for: {file_path}")

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
    return text.strip()

def extract_text_from_docx(file_path: str) -> str:
    return docx2txt.process(file_path).strip()

def extract_text_from_pptx(file_path: str) -> str:
    text = ""
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()